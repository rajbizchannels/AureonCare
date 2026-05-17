from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors (Teal and Golden)
TEAL_PRIMARY = RGBColor(20, 184, 166)    # #14b8a6
TEAL_DARK = RGBColor(15, 118, 110)       # #0f766e
TEAL_LIGHT = RGBColor(94, 234, 212)      # #5eead4
GOLD_PRIMARY = RGBColor(245, 158, 11)    # #f59e0b
GOLD_DARK = RGBColor(217, 119, 6)        # #d97706
GOLD_LIGHT = RGBColor(251, 191, 36)      # #fbbf24
NAVY = RGBColor(30, 41, 59)              # #1e293b
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

def add_title_slide():
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background (simulated with shape)
    background = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = TEAL_DARK
    background.line.fill.background()
    
    # Add logo placeholder text
    logo_box = slide.shapes.add_textbox(
        Inches(3.5), Inches(1),
        Inches(3), Inches(1)
    )
    logo_frame = logo_box.text_frame
    logo_frame.text = "🏥"
    p = logo_frame.paragraphs[0]
    p.font.size = Pt(72)
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "AureonCare"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.5),
        Inches(8), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Executive Presentation - Version 1.2"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = GOLD_LIGHT
    p.alignment = PP_ALIGN.CENTER
    
    # Pre-beta badge
    badge_box = slide.shapes.add_textbox(
        Inches(2.5), Inches(4.5),
        Inches(5), Inches(0.5)
    )
    badge_frame = badge_box.text_frame
    badge_frame.text = "Pre-Beta Development"
    p = badge_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Add shape behind badge
    badge_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(3), Inches(4.5),
        Inches(4), Inches(0.5)
    )
    badge_shape.fill.solid()
    badge_shape.fill.fore_color.rgb = GOLD_PRIMARY
    badge_shape.line.color.rgb = GOLD_DARK
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.5),
        Inches(8), Inches(1)
    )
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Empowering Healthcare with Modern Technology"
    p = tagline_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_items, highlight_text=None):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()
    
    # Top bar
    top_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TEAL_PRIMARY
    top_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEAL_DARK
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        p.level = 0
        p.space_before = Pt(6)
    
    # Highlight box if provided
    if highlight_text:
        highlight_box = slide.shapes.add_shape(
            1, Inches(1), Inches(6), Inches(8), Inches(0.8)
        )
        highlight_box.fill.solid()
        highlight_box.fill.fore_color.rgb = GOLD_LIGHT
        highlight_box.line.color.rgb = GOLD_DARK
        
        text_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(6.1), Inches(7.6), Inches(0.6)
        )
        tf = text_box.text_frame
        tf.text = highlight_text
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY

def add_stats_slide(title, stats_data):
    """Add a slide with statistics cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()
    
    # Top bar
    top_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TEAL_PRIMARY
    top_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEAL_DARK
    
    # Add stat cards in grid
    rows = 2
    cols = 3
    card_width = Inches(2.8)
    card_height = Inches(1.8)
    start_x = Inches(0.7)
    start_y = Inches(2)
    gap_x = Inches(0.2)
    gap_y = Inches(0.3)
    
    for i, (stat_num, stat_label) in enumerate(stats_data[:6]):  # Max 6 stats
        row = i // cols
        col = i % cols
        
        x = start_x + col * (card_width + gap_x)
        y = start_y + row * (card_height + gap_y)
        
        # Card background
        card = slide.shapes.add_shape(1, x, y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = TEAL_PRIMARY
        card.line.color.rgb = TEAL_DARK
        card.line.width = Pt(2)
        
        # Stat number
        num_box = slide.shapes.add_textbox(x, y + Inches(0.3), card_width, Inches(0.8))
        num_frame = num_box.text_frame
        num_frame.text = stat_num
        p = num_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Stat label
        label_box = slide.shapes.add_textbox(x, y + Inches(1.2), card_width, Inches(0.5))
        label_frame = label_box.text_frame
        label_frame.text = stat_label
        label_frame.word_wrap = True
        p = label_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER

def add_table_slide(title, table_data, headers):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()
    
    # Top bar
    top_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TEAL_PRIMARY
    top_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEAL_DARK
    
    # Add table
    rows = len(table_data) + 1  # +1 for header
    cols = len(headers)
    
    left = Inches(0.7)
    top = Inches(1.8)
    width = Inches(8.6)
    height = Inches(4.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    for i in range(cols):
        table.columns[i].width = Inches(width.inches / cols)
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = NAVY
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(241, 245, 249)  # Light gray

# Slide 1: Title
add_title_slide()

# Slide 2: Year 1 Targets
add_stats_slide(
    "🎯 Year 1 Targets & Projections",
    [
        ("5,000+", "Target Users"),
        ("150", "Target Practices"),
        ("2M+", "Platform Capacity"),
        ("14", "Core Modules"),
        ("99.9%", "Uptime SLA"),
        ("4.5★", "Target Satisfaction")
    ]
)

# Slide 3: Key Features
add_content_slide(
    "🚀 Version 1.2 Key Features",
    [
        "🔍 Universal Search - 50-60% projected time savings",
        "📦 Data Archiving - 35-40% DB reduction projected",
        "📋 Audit Logging - Complete HIPAA/SOX compliance",
        "☁️ Cloud Backup - OAuth with Google Drive & OneDrive",
        "📝 SOAP Notes - 20-25% projected denial reduction",
        "🏥 Enhanced Registration - Allergies, PMH, Family History",
        "🔐 Expanded Permissions - Granular RBAC, 14 modules",
        "💡 Help System - AI assistant, 55-60% projected ticket reduction"
    ]
)

# Slide 4: Projected ROI
add_table_slide(
    "💰 Projected First-Year ROI: 650-750%",
    [
        ["Operational Savings", "$450K - $550K"],
        ["• Time savings (50-60%)", "$100K - $120K"],
        ["• Reduced denials (20-25%)", "$160K - $200K"],
        ["Revenue Improvements", "$550K - $665K"],
        ["• Days in A/R (52→35-38)", "$150K - $180K"],
        ["• Collection rate (92%→96-97%)", "$200K - $250K"],
        ["Total Projected Benefit", "$1.0M - $1.2M"],
        ["Implementation Cost", "($150K)"],
        ["Projected Net ROI", "$850K - $1.065M"]
    ],
    ["Category", "Projected Annual Value"]
)

# Slide 5: Development Status
add_stats_slide(
    "🔬 Development Status & Readiness",
    [
        ("95%", "Code Complete"),
        ("100%", "Documentation"),
        ("14", "Modules Ready"),
        ("Q1 2026", "Beta Launch"),
        ("10-15", "Beta Practices"),
        ("Q2 2026", "General Launch")
    ]
)

# Slide 6: Help System
add_content_slide(
    "💡 Help & Documentation System",
    [
        "📚 In-App Help Drawer - Browse, Search, AI Assistant tabs",
        "🤖 AI Assistant - Natural language Q&A, context-aware",
        "📖 12 Comprehensive Guides - Clinical, Billing, Admin",
        "📝 50+ Help Articles - Searchable, categorized, role-based",
        "🌍 Multi-Language - 8 languages supported",
        "",
        "Projected Impact:",
        "• 55-60% reduction in support tickets",
        "• 45-50% faster training for new users"
    ],
    "Complete Billing Documentation: Overview • Payers • Reports • Claims • Payments"
)

# Slide 7: Security & Compliance
add_content_slide(
    "🔒 Security & Compliance",
    [
        "Planned Certifications:",
        "✅ HIPAA Compliant Design",
        "🔄 SOC 2 Type II (Q3 2026)",
        "✅ FDA 21 CFR Part 11",
        "✅ GDPR Ready",
        "✅ 50 State Compliant",
        "",
        "Security Features:",
        "• AES-256 Encryption",
        "• Multi-Factor Auth (MFA)",
        "• Role-Based Access (RBAC)",
        "• 24/7 Monitoring",
        "• Quarterly Pen Testing"
    ],
    "60-70% projected reduction in audit prep time with complete audit trail"
)

# Slide 8: Competitive Advantages
add_table_slide(
    "⚡ Competitive Advantages",
    [
        ["Implementation Time", "30 days (goal)", "90-180 days"],
        ["Universal Search", "14 modules", "2-3 modules"],
        ["AI Help Assistant", "✓ Included", "✗ Not available"],
        ["Languages Supported", "8 languages", "2-3 languages"],
        ["Technology Stack", "Modern (React 18)", "Legacy (2010s)"],
        ["Customer Target", "4.5+/5", "3.8/5 typical"]
    ],
    ["Feature", "AureonCare", "Typical Competitors"]
)

# Slide 9: Development & Launch Roadmap
add_content_slide(
    "🗺️ Development & Launch Roadmap",
    [
        "Q1 2026 (Jan-Mar) - Beta Phase:",
        "🔬 Complete final development & testing",
        "🎯 Launch closed beta (10-15 practices)",
        "🐛 Bug fixes & optimization",
        "✅ Complete security audits",
        "",
        "Q2 2026 (Apr-Jun) - General Launch:",
        "🚀 Official product launch (April 2026)",
        "🎯 Onboard first 50 paying practices",
        "📱 Mobile app (beta) release",
        "🔗 Integration marketplace launch",
        "",
        "Q3-Q4 2026 - Growth & Scale:",
        "🎯 Reach 150 practices by year-end",
        "🤖 AI Clinical Assistant launch",
        "💳 Advanced features rollout"
    ]
)

# Slide 10: Pricing
add_table_slide(
    "💵 Pricing & Early Adopter Benefits",
    [
        ["Essential", "$99", "$950", "2-5 users"],
        ["Professional", "$149", "$1,425", "5-20 users"],
        ["Enterprise", "$199", "$1,900", "20+ users"]
    ],
    ["Plan", "Monthly/User", "Annual/User", "Target Segment"]
)

# Slide 11: Call to Action
add_content_slide(
    "📞 Join the Healthcare Revolution",
    [
        "🎁 Early Adopter Program Benefits:",
        "• 20% discount for first 50 practices",
        "• Priority implementation",
        "• Lifetime VIP support",
        "• Influence product roadmap",
        "",
        "Contact Us Today:",
        "📧 sales@aureoncare.com",
        "📞 1-800-AUREON1",
        "🌐 https://aureoncare.com",
        "",
        "For Healthcare Practices:",
        "Schedule demo • Join beta • Request pricing • Reserve spot",
        "",
        "For Investors & Partners:",
        "Investment overview • Financial projections • Partnership opportunities"
    ],
    "Beta: Q1 2026 | Launch: Q2 2026 | Beta & Early Adopter Slots Available"
)

# Slide 12: Summary
add_stats_slide(
    "The Future of Healthcare Starts Here",
    [
        ("650-750%", "Projected ROI"),
        ("30-Day", "Implementation"),
        ("14", "Integrated Modules"),
        ("Q1→Q2", "Beta→Launch"),
        ("$850K+", "Net Benefit"),
        ("20%", "Early Adopter Discount")
    ]
)

# Save presentation
prs.save('/home/user/MedFlow/AureonCare_Executive_Presentation.pptx')
print("✅ PowerPoint presentation created successfully!")
print("📁 File: /home/user/MedFlow/AureonCare_Executive_Presentation.pptx")
